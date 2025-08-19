import streamlit as st
import re
import json
from io import BytesIO
from PIL import Image
import pytesseract
import docx2txt
import fitz # PyMuPDF
from fpdf import FPDF
import pandas as pd

# Optional imports for document types
try:
    from docx import Document
    from docx.shared import Inches
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    st.warning("`python-docx` not installed. DOCX file processing may be limited to basic text extraction and DOCX generation will not be available.")

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    PPTS_AVAILABLE = True
except ImportError:
    PPTS_AVAILABLE = False
    st.warning("`python-pptx` not installed. PPTX generation will not be available.")

# --- TEXT EXTRACTION FUNCTION ---
def extract_text(file):
    text = ""
    try:
        if file.type == "application/pdf":
            with fitz.open(stream=file.read(), filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = docx2txt.process(file)
        elif file.type.startswith("text"):
            text = file.read().decode("utf-8", errors="ignore")
        elif file.type.startswith("image"):
            img = Image.open(file)
            text = pytesseract.image_to_string(img)
        else:
            st.error(f"Unsupported file type: {file.type}. Please upload PDF, DOCX, TXT, or image files.")
            return ""
    except Exception as e:
        st.error(f"Error extracting text from {file.name}: {e}. Ensure all dependencies (like Tesseract for images) are correctly set up.")
        return ""
    return text

# --- FIELD EXTRACTION FUNCTION ---
def extract_field(text, patterns):
    """
    Extracts a field value from text using a list of regex patterns.
    Returns the first clean match it finds.
    """
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
        if match:
            if match.groups():
                for group in match.groups():
                    if group:
                        clean_value = re.sub(r"\s+", " ", group.strip())
                        # FIX: Corrected the word boundary characters from '' to '\b'
                        clean_value = re.split(
                            r"\b(Branch|Date|Note|Code|Signature|Time|No\.|Number|Claim|Policy|Name|Gender|Age|Address|Email|Mobile|Contact|Phone|Relationship|Occupation)\b",
                            clean_value,
                            maxsplit=1
                        )[0].strip()
                        if clean_value.lower() in [
                            "of", "the", "of the", "name", "aadhar holder", "n", "claim", "-", "not mentioned",
                            "rs", "no", "number", "date", "place", "contact", "address",
                            "gender", "age", "mobile", "email", "phone", "details", "birth", "ion", "occu", "na", "nil"
                        ] or len(clean_value.strip()) < 3:
                            continue
                        if clean_value.lower() == "na":
                            return "NA"
                        if len(clean_value) < 5 and not any(char.isdigit() for char in clean_value) and not any(char.isalpha() for char in clean_value):
                            continue
                        return clean_value
            # If no groups, return the full match but cleaned
            return re.sub(r"\s+", " ", match.group(0).strip())
    return ""

def extract_and_format_raw_block(field_value):
    """
    Parses and formats raw text block with descriptions and amounts.
    Handles hospital bills, repair costs, pharmacy medication tables, and fallbacks.
    """
    if not field_value:
        return ""

    # --- 🛠️ Preprocess: Fix line breaks for multiline descriptions ---
    lines = field_value.splitlines()
    fixed_lines = []
    buffer = ""

    for line in lines:
        line = line.strip()
        if not line:
            continue

        if buffer:
            if re.match(r".*[\(\,:\-]$", buffer) or len(buffer.split()) < 3:
                buffer += " " + line
            else:
                fixed_lines.append(buffer)
                buffer = line
        else:
            buffer = line
    if buffer:
        fixed_lines.append(buffer)

    field_value = "\n".join(fixed_lines)

    field_value = field_value.replace("Replace- ment", "Replacement")
    field_value = re.sub(r'"\s*\n\s*"', ' ', field_value)
    field_value = re.sub(r'\s{2,}', ' ', field_value).strip()

    formatted = ""
    total = 0

    # --- Pattern 1: Hospital-style bill breakup ---
    pattern_bill_items = r"(\d+)\s*.*?\s*\"([^\"]*?)\".*?([\d,]+(?:\.\d{1,2})?)\b"
    matches_bill_items = re.findall(pattern_bill_items, field_value, re.IGNORECASE | re.DOTALL)
    if matches_bill_items:
        for s_no, desc, amt in matches_bill_items:
            desc = desc.strip()
            amt_clean = amt.replace(",", "")
            if amt_clean.replace('.', '').isdigit():
                amount_num = float(amt_clean)
                formatted += f"- {desc}: Rs. {amount_num:,.0f}\n"
                total += amount_num
        if total > 0:
            formatted += f"- Total: Rs. {total:,.0f}"
        return formatted

    # --- Pattern 2: Estimate-style (desc + 2+ spaces + amount) ---
    pattern_estimate_lines = r"([A-Za-z0-9(),\s\-\/\\.&]+?)\s{2,}(Rs\.?\s*|\₹?)\s*(\d{1,3}(?:,\d{3})*(?:\.\d{1,2})?)\b"
    matches_estimate = re.findall(pattern_estimate_lines, field_value, re.IGNORECASE)
    if matches_estimate:
        for desc, _, amt in matches_estimate:
            desc = desc.strip()
            amt_clean = amt.replace(",", "")
            if amt_clean.replace('.', '').isdigit():
                amount_num = float(amt_clean)
                formatted += f"- {desc}: Rs. {amount_num:,.0f}\n"
                total += amount_num
        if total > 0:
            formatted += f"- Total: Rs. {total:,.0f}"
        return formatted

    # --- Pattern 3: Motor Repair-style items ---
    field_value = re.sub(r"Total Repair Cost[:\s₹Rs\.]*[\d,]+", "", field_value, flags=re.IGNORECASE)
    pattern_motor_repair = r"(\d+)\s+([A-Za-z0-9(),\s\-\/\\.&]+?)\s+(\d+)\s+([\d,]+)"
    matches_motor_repair = re.findall(pattern_motor_repair, field_value)
    if matches_motor_repair:
        for num, desc, qty, amt in matches_motor_repair:
            desc = desc.strip()
            amt_clean = amt.replace(",", "")
            if amt_clean.replace('.', '').isdigit():
                amount_num = float(amt_clean)
                formatted += f"{num}. {desc} - Qty: {qty}, Cost: Rs. {amount_num:,.0f}\n"
                total += amount_num
        if total > 0:
            formatted += f"\nTotal Repair Cost: Rs. {total:,.0f}"
        return formatted

    # --- NEW Pattern: Pharmacy (S. No. + Med Name + Final Amount) ---
    pattern_pharmacy_sn_amt = r"(\d+)\s+(Tab\.|Cap\.|Syp\.|Inj\.|Oint\.|Cream\.)\s+([A-Za-z0-9\s\-\.]+?)\s+\d+\s+\d+\s+([\d,]+(?:\.\d{1,2})?)"
    matches_pharmacy_sn_amt = re.findall(pattern_pharmacy_sn_amt, field_value)
    if matches_pharmacy_sn_amt:
        for s_no, form, name, amount in matches_pharmacy_sn_amt:
            desc = f"{form} {name.strip()}"
            amt_clean = amount.replace(",", "")
            if amt_clean.replace('.', '').isdigit():
                amount_num = float(amt_clean)
                formatted += f"{s_no}. {desc} - Rs. {amount_num:,.0f}\n"
                total += amount_num
        if total > 0:
            formatted += f"Total: Rs. {total:,.0f}"
        return formatted
    # --- Pattern 5: Fallback "X: Rs 123" ---
    pattern_simple_desc_amount = r"([A-Za-z0-9\s,()\/\-]+?)\s*:\s*Rs\.?\s*([\d,]+(?:\.\d{1,2})?)"
    matches_simple = re.findall(pattern_simple_desc_amount, field_value)
    if matches_simple:
        for desc, amt in matches_simple:
            desc = desc.strip()
            amt_clean = amt.replace(",", "")
            if amt_clean.replace('.', '').isdigit():
                amount_num = float(amt_clean)
                formatted += f"- {desc}: Rs. {amount_num:,.0f}\n"
                total += amount_num
        if total > 0:
            formatted += f"- Total: Rs. {total:,.0f}"
        return formatted

    # --- Pattern 6: Fallback meds like "Tab. XYZ: Rs 123" ---
    pattern_fallback_meds = r"(Tab\.|Cap\.|Syp\.|Inj\.|Oint\.|Cream\.)\s+([A-Za-z0-9\s\-\.]+):\s*Rs\.?\s*([\d,]+(?:\.\d{1,2})?)"
    matches_fallback_meds = re.findall(pattern_fallback_meds, field_value)
    if matches_fallback_meds:
        for form, name, amt in matches_fallback_meds:
            desc = f"{form} {name.strip()}"
            amt_clean = amt.replace(",", "")
            if amt_clean.replace('.', '').isdigit():
                amount_num = float(amt_clean)
                formatted += f"• {desc} - ₹{amount_num:,.0f}\n"
                total += amount_num
        if total > 0:
            formatted += f"Total: ₹{total:,.0f}"
        return formatted

    # --- Final Fallback: Raw text ---
    if not formatted:
        clean_raw_value = field_value.strip().replace("\n", " ").replace("\r", " ")
        if clean_raw_value:
            return "- " + clean_raw_value

    return formatted
# --- SUMMARY EXTRACTION ---
def extract_summary(text, config):
    """
    Extracts summary fields based on the provided configuration.
    """
    summary = {}
    for field_name, patterns in config.get("fields", {}).items():
        value = extract_field(text, patterns)

        if value.strip():  # Only add if a non-empty value was found
            # Fix hyphenated line breaks like 'Rear-end col-\nlision' → 'Rear-end collision'
            value = re.sub(r'(\w)-\s+(\w)', r'\1\2', value)

            # Final cleanup for trailing junk like "Nationality", "By submitting..."
            cleaned_value = re.sub(
                r"\b(Nationality|Claimant|SPAARC.*|By submitting.*|NEFT mandate.*|insurance claim.*)$",
                "", value, flags=re.IGNORECASE).strip()

            # Ensure Rs. prefix for specific monetary fields
            money_fields = [
                "estimated cost of repairs", "total repair cost", "repair estimate",
                "amount claimed", "claim amount", "sum insured", "approved amount",
                "depreciation amount"
            ]

            if field_name.lower() in money_fields:
                # Add Rs. if missing and value looks like a number
                if re.match(r"^\d{1,3}(,\d{3})*(\.\d{1,2})?$", cleaned_value):
                    cleaned_value = f"Rs. {cleaned_value}"
                elif re.match(r"^\d+(\.\d+)?$", cleaned_value):
                    cleaned_value = f"Rs. {int(float(cleaned_value)):,.0f}"

            summary[field_name] = cleaned_value

    return summary

# --- PDF GENERATION ---
# No custom header for logo for now
class MyFPDF(FPDF):
    def header(self):
        # Header is empty for now as requested
        pass

def generate_pdf(text_content):
    """Generates a PDF from a given text string with a blank second page."""
    pdf = MyFPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=10) # Smaller font for better fit

    # Split text content to place some on the first page, and the rest later
    # This is a simple split, you might need more sophisticated logic
    lines = text_content.splitlines()
    first_page_content_lines = lines[:10] # Example: Put first 10 lines on page 1
    remaining_content_lines = lines[10:]

    pdf.multi_cell(0, 5, "\n".join(first_page_content_lines).encode('latin-1', 'replace').decode('latin-1'))

    # Add a blank second page
    pdf.add_page() # This adds a new page, which will be blank
    # No content added to this page, making it blank

    # Add remaining content to subsequent pages
    if remaining_content_lines:
        pdf.add_page() # Add a new page for the rest of the content
        pdf.multi_cell(0, 5, "\n".join(remaining_content_lines).encode('latin-1', 'replace').decode('latin-1'))

    return pdf.output(dest='S').encode('latin-1')

# --- DOCX GENERATION ---
def generate_docx(text_content):
    """Generates a DOCX from a given text string with a blank second page."""
    if not DOCX_AVAILABLE:
        st.error("`python-docx` is not installed. Cannot generate DOCX file.")
        return None
    document = Document()

    # No logo added to header for now

    document.add_paragraph(text_content) # Add all content

    # Add a page break to create the blank second page
    document.add_page_break()
    document.add_paragraph("") # Add an empty paragraph to ensure a new, blank page is created
    document.add_page_break() # Add another page break for content after the blank page

    # Save to a BytesIO object
    buffer = BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

# --- PPTX GENERATION ---
def generate_pptx(text_content):
    """Generates a simple PPTX from a given text string."""
    if not PPTS_AVAILABLE:
        st.error("`python-pptx` is not installed. Cannot generate PPTX file.")
        return None

    prs = Presentation()
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    # body = slide.placeholders[1] # For subtitle, if desired

    title.text = "Insurance Claim Summary"
    # body.text = "Extracted Information"

    # Add content slides for the formatted text
    bullet_slide_layout = prs.slide_layouts[1] # Title and Content layout

    # Split the formatted text into manageable chunks for slides
    # A simple approach: split by sections or a fixed number of lines
    sections = text_content.split('### ') # Assuming '###' marks new sections

    for i, section_text in enumerate(sections):
        if not section_text.strip():
            continue

        # The first part might be empty or a general intro if it was before the first '###'
        if i == 0 and not section_text.strip().splitlines()[0].strip():
            lines = section_text.strip().splitlines()
            if len(lines) > 0:
                slide_title = lines[0].strip() # Use first line as title
                slide_content = "\n".join(lines[1:]).strip()
            else:
                slide_title = "Summary Details"
                slide_content = section_text.strip()
        else:
            lines = section_text.strip().splitlines()
            slide_title = lines[0].strip() if lines else "Details"
            slide_content = "\n".join(lines[1:]).strip()

        # Create a new slide for each section
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = slide_title

        # Add content as bullet points or just text
        tf = body.text_frame
        tf.clear() # Clear existing text in placeholder
        p = tf.add_paragraph()
        p.text = slide_content
        # You might want to split slide_content further into bullet points if needed

    # Save to a BytesIO object
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# --- STREAMLIT UI ---
st.set_page_config(layout="wide", page_title="Insurance Claim Summarizer")
st.sidebar.title("Insurance Claim Document Summarizer")
st.sidebar.markdown("Upload your insurance claim documents (PDF, DOCX, TXT, Images) to extract key information and generate a structured summary.")

# Original file uploader (not a radio button for upload)
st.sidebar.header("1. Upload Document(s)")
uploaded_files = st.sidebar.file_uploader(
    "Choose file(s)",
    type=["pdf", "png", "jpg", "jpeg", "txt", "docx"],
    accept_multiple_files=True
)

st.sidebar.header("2. Select Document Type")
# Changed from selectbox to radio button
document_type = st.sidebar.radio(
    "Document Type",
    ["Vehicle Insurance", "Health Insurance", "Life Insurance"], # "Motor Insurance" to "Vehicle Insurance"
    key="document_type_radio", # Added a key for radio button
    index=0 # Default to Vehicle Insurance
)

st.title("📋 Insurance Claim Summary Generator")
st.markdown("---")

# Initialize combined_summary to ensure it's always available
combined_summary = {}
all_file_texts = {} # Store raw text for display

if uploaded_files: # Process general uploaded files first
    st.subheader("📄 Extracted Text Previews")


    # Load configuration once per document type
    config_path = f"configs/{document_type.lower().replace(' ', '_')}_config.json"
    config = {"fields": {}, "summary_sections": []} # Default empty config
    try:
        with open(config_path, "r", encoding="utf-8") as f:
            config = json.load(f)
        st.sidebar.success(f"Loaded configuration for '{document_type}'.")
    except FileNotFoundError:
        st.error(f"Configuration file not found for {document_type}. Please ensure '{config_path}' exists.")
        st.info("Example: Create a `configs` folder and add `vehicle_insurance_config.json` inside it.") # Updated config name hint
    except json.JSONDecodeError:
        st.error(f"Error decoding JSON from {config_path}. Please check the file's format (e.g., missing commas, extra commas, unclosed brackets).")

    # Process each uploaded file from general uploader
    for i, uploaded_file in enumerate(uploaded_files):
        st.markdown(f"**File {i+1}: {uploaded_file.name}**")
        with st.spinner(f"Extracting text from {uploaded_file.name}..."):
            file_text = extract_text(uploaded_file)

        if file_text:
            all_file_texts[uploaded_file.name] = file_text
            st.expander(f"Click to view raw text for {uploaded_file.name}").text(file_text[:1500] + ("..." if len(file_text) > 1500 else ""))

            # Extract summary for the current file
            file_summary = extract_summary(file_text, config)

            # Merge into combined_summary, prioritizing the first non-empty value found
            for key, value in file_summary.items():
                if key not in combined_summary or not combined_summary[key].strip():
                    combined_summary[key] = value

    if combined_summary:
        # --- ✅ Display structured summary ---
        st.subheader("✅ Final Structured Summary (from all documents)")

        # KYC and Policy Assignment radio buttons (instead of dropdowns)
        policy_assigned = st.radio(
            "Could you please confirm whether the policy has been assigned?",
            ("No", "Yes"),
            key="policy_assigned_radio", # Changed key from select to radio
            index=0 # Default to "No"
        )
        kyc_verified = st.radio(
            "Could you please let us know whether the KYC verification has been completed?",
            ("Yes", "No"), # Default to Yes
            key="kyc_verified_radio", # Changed key from select to radio
            index=0
        )

        if kyc_verified == "No":
            st.warning("Please verify the KYC.")

        st.markdown("---") # Add a separator for clarity after these statements

        # --- Health Insurance specific questions placed here ---
        if document_type == "Health Insurance":
            st.subheader("Hospital & Claim Processing Information")

            hospital_type_network = st.radio(
                "Is the hospital type network?",
                ("No", "Yes"),
                key="hospital_type_network_radio"
            )
            combined_summary["Hospital Type Network"] = hospital_type_network # Store for download

            if hospital_type_network == "Yes":
                cashless_facility = st.radio(
                    "Do you wish to avail the cashless facility?",
                    ("No", "Yes"),
                    key="cashless_facility_radio"
                )
                combined_summary["Cashless Facility Availed"] = cashless_facility # Store for download

                if cashless_facility == "Yes":
                    # Changed to accept multiple files
                    health_card_files = st.file_uploader(
                        "Upload Health Insurance Card(s)",
                        type=["pdf", "png", "jpg", "jpeg", "txt"],
                        accept_multiple_files=True, # ALLOW MULTIPLE FILES
                        key="health_card_upload"
                    )
                    if health_card_files:
                        for idx, h_file in enumerate(health_card_files):
                            st.info(f"Health Insurance Card '{h_file.name}' uploaded.")
                            # Optional: extract text from health card if needed
                            # health_card_text = extract_text(h_file)
                            # if health_card_text:
                            #     all_file_texts[f"HealthCard_{h_file.name}"] = health_card_text
                            #     st.expander(f"Raw text for {h_file.name}").text(health_card_text[:500] + "...")
                else: # Cashless facility No
                    # Changed to accept multiple files
                    hospital_bills_files = st.file_uploader(
                        "Upload Hospital Bill(s)",
                        type=["pdf", "png", "jpg", "jpeg", "txt", "docx"],
                        accept_multiple_files=True, # ALLOW MULTIPLE FILES
                        key="hospital_bills_upload_cashless_no" # Unique key
                    )
                    if hospital_bills_files:
                        for idx, bill_file in enumerate(hospital_bills_files):
                            st.info(f"Hospital Bill '{bill_file.name}' uploaded.")
                            bill_text = extract_text(bill_file)
                            if bill_text:
                                all_file_texts[f"HospitalBill_{bill_file.name}"] = bill_text
                                st.expander(f"Click to view raw text for {bill_file.name}").text(bill_text[:1500] + ("..." if len(bill_text) > 1500 else ""))
                                bill_summary = extract_summary(bill_text, config)
                                for key, value in bill_summary.items():
                                    if key not in combined_summary or not combined_summary[key].strip():
                                        combined_summary[key] = value
            else: # Hospital type Network No - THIS IS THE CHANGE BLOCK
                st.info("It is a reimbursement claim.")
                combined_summary["Claim Type"] = "Reimbursement Claim" # Store for download

                # Allow uploading bills for reimbursement claim as well
                hospital_bills_files_reimbursement = st.file_uploader(
                    "Upload Hospital Bill(s) for Reimbursement", # Changed label for clarity
                    type=["pdf", "png", "jpg", "jpeg", "txt", "docx"],
                    accept_multiple_files=True,
                    key="hospital_bills_upload_reimbursement" # Unique key
                )
                if hospital_bills_files_reimbursement:
                    for idx, bill_file_reimbursement in enumerate(hospital_bills_files_reimbursement):
                        st.info(f"Hospital Bill '{bill_file_reimbursement.name}' uploaded for reimbursement.")
                        bill_text_reimbursement = extract_text(bill_file_reimbursement)
                        if bill_text_reimbursement:
                            all_file_texts[f"ReimbursementBill_{bill_file_reimbursement.name}"] = bill_text_reimbursement
                            st.expander(f"Click to view raw text for {bill_file_reimbursement.name}").text(bill_text_reimbursement[:1500] + ("..." if len(bill_file_reimbursement.name) > 1500 else ""))
                            bill_summary_reimbursement = extract_summary(bill_text_reimbursement, config)
                            for key, value in bill_summary_reimbursement.items():
                                if key not in combined_summary or not combined_summary[key].strip():
                                    combined_summary[key] = value

            st.markdown("---") # Another separator after health insurance questions


        # --- Existing Motor Insurance FIR Questions ---
        if document_type == "Vehicle Insurance": # Changed to Vehicle Insurance
            st.subheader("Police & Authority Information")
            fir_status_selection = st.radio(
                "Has FIR/affidavit been filed?",
                ("No", "Yes"),
                key="fir_status_radio_main_page"
            )
            fir_affidavit_files_in_page = None # Changed variable name to plural
            if fir_status_selection == "Yes":
                # Changed to accept multiple files
                fir_affidavit_files_in_page = st.file_uploader(
                    "Upload FIR/Affidavit Document(s)",
                    type=["pdf", "png", "jpg", "jpeg", "txt", "docx"],
                    accept_multiple_files=True, # ALLOW MULTIPLE FILES
                    key="fir_affidavit_upload"
                )
            elif fir_status_selection == "No":
                st.info("No supporting police document is required.")

            if fir_affidavit_files_in_page: # Check if files were actually uploaded
                for idx, fir_file in enumerate(fir_affidavit_files_in_page): # Loop through files
                    st.info(f"{fir_status_selection} document '{fir_file.name}' uploaded.")
                    with st.spinner(f"Extracting text from {fir_file.name}..."):
                        fir_affidavit_text = extract_text(fir_file)
                    if fir_affidavit_text:
                        all_file_texts[f"FIR_Affidavit_{fir_file.name}"] = fir_affidavit_text
                        st.expander(f"Click to view raw text for {fir_file.name}").text(fir_affidavit_text[:1500] + ("..." if len(fir_affidavit_text) > 1500 else ""))

                        fir_affidavit_summary = extract_summary(fir_affidavit_text, config)
                        for key, value in fir_affidavit_summary.items():
                            if key not in combined_summary or not combined_summary[key].strip():
                                combined_summary[key] = value
            combined_summary["FIR Status"] = fir_status_selection
            st.markdown("---") # Another separator after motor insurance questions


        # --- Building the formatted summary string for display ---
        formatted = ""
        for section in config.get("summary_sections", []):
            heading = section.get("section_title")
            section_fields = section.get("fields", [])
            section_content = []

            for field in section_fields:
                value = combined_summary.get(field, "").strip()
                if value:
                    short_field = field # Initialize with original

                    # --- Field name shortening logic ---
                    if heading and heading.lower() == "life assured details" and field.lower().startswith("insured "):
                        short_field = field[len("Insured "):].strip(" :-")
                    elif heading and heading.lower() == "bank & payout details" and field.lower().startswith("claimant "):
                        short_field = field[len("Claimant "):].strip(" :-")
                    elif heading and heading.lower() == "claim submission details" and (" (official)" in field.lower() or " (official use)" in field.lower()):
                        short_field = field.replace(" (Official)", "").replace(" (Official Use)", "").strip(" :-")
                    elif heading and heading.lower() == "death certificate details" and " (death certificate)" in field.lower():
                        short_field = field.replace("(Death Certificate)", "").strip(" :-")
                    elif heading and heading.lower() == "kyc details" and " (kyc)" in field.lower():
                        short_field = field.replace("(KYC)", "").strip(" :-")
                    elif heading and heading.lower() == "policy details" and field.lower() == "policy number":
                        short_field = "Policy No."
                    elif heading and heading.lower() == "claimant details" and field.lower().startswith("claimant "):
                        short_field = field[len("Claimant "):].strip(" :-")
                    elif heading and heading.lower() == "driver details" and field.lower().startswith("driver "):
                        if "driving" not in field.lower():
                            short_field = field[len("Driver "):].strip(" :-")
                        else:
                            short_field = field
                    elif heading and heading.lower() == "garage details" and field.lower().startswith("garage "):
                        short_field = field[len("Garage "):].strip(" :-")
                    elif heading and heading.lower() == "other insurance details" and field.lower().startswith("other insurance - "):
                        short_field = field[len("Other Insurance - "):].strip(" :-")
                    elif heading and heading.lower() == "interest holder details" and field.lower().startswith("interest holder - "):
                        short_field = field[len("Interest Holder - "):].strip(" :-")
                    elif heading and heading.lower() == "discharge voucher" and "(discharge voucher)" in field.lower():
                        short_field = field.replace("(Discharge Voucher)", "").strip(" :-")
                    elif heading and heading.lower() == "satisfaction note" and "(satisfaction note)" in field.lower():
                        short_field = field.replace("(Satisfaction Note)", "").strip(" :-")
                    elif heading and heading.lower() == "kyc information" and "(kyc)" in field.lower():
                        short_field = field.replace("(KYC)", "").strip(" :-")
                    elif heading and heading.lower() == "primary insured details" and field.lower().startswith("primary insured "):
                        short_field = field[len("Primary Insured "):].strip(" :-")
                    elif heading and heading.lower() == "insured person hospitalized details" and field.lower().startswith("insured person hospitalized "):
                        short_field = field[len("Insured Person Hospitalized "):].strip(" :-")
                    elif heading and field.lower().startswith(heading.lower().replace(" details", "").strip() + " "):
                        short_field = field[len(heading.lower().replace(" details", "").strip()):].strip(" :-")
                    elif field in ["FIR Status", "Hospital Type Network", "Cashless Facility Availed", "Claim Type"]: # Exclude new interactive fields
                        continue # Skip adding these to the formatted string directly from combined_summary here

                    # Special formatting for raw blocks
                    if any(raw_key in field for raw_key in [
                        "Bill Breakup Details Raw", "Medication Details Raw",
                        "Repair Items Raw", "Loss Details Raw", "Cost Details Raw"
                    ]):
                        section_content.append(f"### {short_field}\n" + extract_and_format_raw_block(value))
                    else:
                        section_content.append(f"- {short_field}: {value}")

            # Add section to formatted string if it has content
            if section_content:
                formatted += f"\n### {heading}\n" + "\n".join(section_content) + "\n"

        st.markdown(formatted) # Display the rest of the formatted summary

        # --- 🔽 Download Section ---
        st.subheader("🔽 Download Summary")

        # Changed from selectbox to radio buttons
        download_format = st.radio(
            "Select download format:",
            ["PDF", "Word (DOCX)", "PowerPoint (PPTX)"],
            key="download_format_radio", # Changed key to reflect radio
            index=0 # Default to PDF
        )

        download_data = None
        download_filename = "insurance_summary"
        download_mime = "application/octet-stream"

        # Prepare full content for downloads (PDF, DOCX, PPTX)
        initial_statements_for_download = f"Could you please confirm whether the policy has been assigned?: {policy_assigned}\nCould you please let us know whether the KYC verification has been completed?: {kyc_verified}\n"

        # Add FIR status to download content if applicable (it's already in combined_summary now)
        if "FIR Status" in combined_summary:
            initial_statements_for_download += f"FIR/Affidavit Status: {combined_summary['FIR Status']}\n"

        # Add Health Insurance specific questions to download content
        if document_type == "Health Insurance":
            if "Hospital Type Network" in combined_summary:
                initial_statements_for_download += f"Is the hospital type network?: {combined_summary['Hospital Type Network']}\n"
            if "Cashless Facility Availed" in combined_summary:
                initial_statements_for_download += f"Do you wish to avail the cashless facility?: {combined_summary['Cashless Facility Availed']}\n"
            if "Claim Type" in combined_summary:
                initial_statements_for_download += f"Claim Type: {combined_summary['Claim Type']}\n"


        full_text_for_download = initial_statements_for_download + "\n" + formatted


        if download_format == "PDF":
            download_data = generate_pdf(full_text_for_download)
            download_filename += ".pdf"
            download_mime = "application/pdf"
        elif download_format == "Word (DOCX)":
            if DOCX_AVAILABLE:
                download_data = generate_docx(full_text_for_download)
                download_filename += ".docx"
                download_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            else:
                st.warning("Cannot generate Word (DOCX) file. `python-docx` library is not installed.")
        elif download_format == "PowerPoint (PPTX)":
            if PPTS_AVAILABLE:
                download_data = generate_pptx(full_text_for_download)
                download_filename += ".pptx"
                download_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                st.warning("Cannot generate PowerPoint (PPTX) file. `python-pptx` library is not installed.")

        if download_data:
            st.download_button(
                f"Download as {download_format}",
                data=download_data,
                file_name=download_filename,
                mime=download_mime,
                help=f"Download the structured summary as a {download_format} file."
            )
        else:
            st.warning("No relevant information could be extracted from the uploaded files based on the selected document type configuration. Please check the file content or try a different document type.")
else:
    st.info("📁 Upload one or more documents from the left sidebar to begin. Select the document type to ensure accurate extraction.")
    st.markdown("---")
    st.markdown("### How to use:")
    st.markdown("1.  **Upload:** Drag and drop your insurance documents (PDF, DOCX, TXT, JPG, PNG) into the sidebar.")
    st.markdown("2.  **Select Type:** Choose the relevant insurance type (Vehicle, Health, Life) using the radio buttons.")
    st.markdown("3.  **View & Download:** The extracted summary will appear below. Select your desired download format using the radio buttons and click 'Download'.")     